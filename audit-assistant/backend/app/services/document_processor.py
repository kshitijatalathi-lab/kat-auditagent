"""Document processing service for handling various file formats."""
import io
import logging
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Union

import docx
import pandas as pd
import pdfplumber
import pytesseract
from PIL import Image
from pptx import Presentation

from app.core.config import settings
from app.utils.file_utils import get_content_type, is_file_type_supported

logger = logging.getLogger(__name__)

# Configure Tesseract path if specified in settings
if settings.TESSERACT_CMD:
    pytesseract.pytesseract.tesseract_cmd = settings.TESSERACT_CMD

class DocumentProcessor:
    """Process various document formats to extract text content."""
    
    @staticmethod
    def process_file(file_path: Union[str, Path]) -> Dict[str, str]:
        """
        Process a file and extract its text content.
        
        Args:
            file_path: Path to the file to process
            
        Returns:
            Dictionary containing the extracted text and metadata
        """
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # Get file extension and content type
        ext = file_path.suffix.lower()
        content_type = get_content_type(file_path.name)
        
        # Process based on file type
        if ext == '.pdf':
            text = DocumentProcessor._extract_pdf_text(file_path)
        elif ext in ['.docx', '.doc']:
            text = DocumentProcessor._extract_docx_text(file_path)
        elif ext == '.txt':
            text = DocumentProcessor._extract_txt_text(file_path)
        elif ext in ['.xlsx', '.xls', '.ods']:
            text = DocumentProcessor._extract_excel_text(file_path)
        elif ext in ['.pptx', '.ppt', '.odp']:
            text = DocumentProcessor._extract_pptx_text(file_path)
        elif ext in ['.jpg', '.jpeg', '.png', '.tiff', '.bmp']:
            text = DocumentProcessor._extract_image_text(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")
        
        return {
            'content': text,
            'content_type': content_type,
            'file_size': file_path.stat().st_size,
            'page_count': DocumentProcessor._get_page_count(file_path, ext),
        }
    
    @staticmethod
    def _extract_pdf_text(file_path: Path) -> str:
        """Extract text from PDF file."""
        text = []
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text.append(page.extract_text() or "")
            return "\n\n".join(filter(None, text))
        except Exception as e:
            logger.error(f"Error extracting text from PDF {file_path}: {e}")
            return ""
    
    @staticmethod
    def _extract_docx_text(file_path: Path) -> str:
        """Extract text from DOCX file."""
        try:
            doc = docx.Document(file_path)
            return "\n\n".join([paragraph.text for paragraph in doc.paragraphs])
        except Exception as e:
            logger.error(f"Error extracting text from DOCX {file_path}: {e}")
            return ""
    
    @staticmethod
    def _extract_txt_text(file_path: Path) -> str:
        """Extract text from plain text file."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Error reading text file {file_path}: {e}")
            return ""
    
    @staticmethod
    def _extract_excel_text(file_path: Path) -> str:
        """Extract text from Excel file."""
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            text_parts = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                text_parts.append(f"--- Sheet: {sheet_name} ---\n{df.to_string()}")
            
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting text from Excel {file_path}: {e}")
            return ""
    
    @staticmethod
    def _extract_pptx_text(file_path: Path) -> str:
        """Extract text from PowerPoint file."""
        try:
            prs = Presentation(file_path)
            text = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                if slide_text:
                    text.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))
            
            return "\n\n".join(text)
        except Exception as e:
            logger.error(f"Error extracting text from PowerPoint {file_path}: {e}")
            return ""
    
    @staticmethod
    def _extract_image_text(file_path: Path) -> str:
        """Extract text from image using OCR."""
        try:
            # Open the image file
            image = Image.open(file_path)
            
            # Convert to grayscale for better OCR
            if image.mode != 'L':
                image = image.convert('L')
            
            # Use Tesseract to do OCR on the image
            return pytesseract.image_to_string(image)
        except Exception as e:
            logger.error(f"Error performing OCR on image {file_path}: {e}")
            return ""
    
    @staticmethod
    def _get_page_count(file_path: Path, file_ext: str) -> int:
        """Get the number of pages in a document."""
        try:
            if file_ext == '.pdf':
                with pdfplumber.open(file_path) as pdf:
                    return len(pdf.pages)
            elif file_ext in ['.docx', '.doc']:
                doc = docx.Document(file_path)
                return len(doc.paragraphs)  # Rough estimate
            elif file_ext in ['.pptx', '.ppt', '.odp']:
                prs = Presentation(file_path)
                return len(prs.slides)
            elif file_ext in ['.xlsx', '.xls', '.ods']:
                return len(pd.ExcelFile(file_path).sheet_names)
            else:
                return 1
        except Exception as e:
            logger.error(f"Error getting page count for {file_path}: {e}")
            return 1
